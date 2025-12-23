from pptlint.utils import iter_text_shapes
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tests.fakes import FakeShape, FakeSlide, FakeTextFrame


def test_iter_text_shapes_yields_only_text_shapes():
    text_shape = FakeShape(text_frame=FakeTextFrame(text="Hello"))
    non_text_shape = FakeShape(text_frame=None)
    slide = FakeSlide([text_shape, non_text_shape])

    shapes = list(iter_text_shapes(slide))

    assert shapes == [text_shape]


def test_iter_text_shapes_walks_group_shapes():
    inner_text = FakeShape(text_frame=FakeTextFrame(text="Inner"))
    inner_non_text = FakeShape(text_frame=None)
    group = FakeShape(shapes=[inner_text, inner_non_text], shape_type=MSO_SHAPE_TYPE.GROUP)
    outer_text = FakeShape(text_frame=FakeTextFrame(text="Outer"))
    slide = FakeSlide([group, outer_text])

    shapes = list(iter_text_shapes(slide))

    assert inner_text in shapes
    assert outer_text in shapes
    assert inner_non_text not in shapes
