import sys
import types
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))


try:
    import pptx
except Exception:
    pptx = types.ModuleType("pptx")
    pptx.enum = types.ModuleType("pptx.enum")
    pptx.enum.shapes = types.ModuleType("pptx.enum.shapes")
    pptx.oxml = types.ModuleType("pptx.oxml")
    pptx.oxml.ns = types.ModuleType("pptx.oxml.ns")
    pptx.util = types.ModuleType("pptx.util")
    pptx.exc = types.ModuleType("pptx.exc")

    class _MSO_SHAPE_TYPE:
        GROUP = 6

    class _PP_PLACEHOLDER:
        TITLE = 1
        CENTER_TITLE = 2

    class _PackageNotFoundError(Exception):
        pass

    def _qn(tag):
        return tag

    class _Length:
        def __init__(self, pt):
            self.pt = pt

    def _Presentation(_):
        return None

    def _Pt(val):
        return _Length(val)

    pptx.enum.shapes.MSO_SHAPE_TYPE = _MSO_SHAPE_TYPE
    pptx.enum.shapes.PP_PLACEHOLDER = _PP_PLACEHOLDER
    pptx.oxml.ns.qn = _qn
    pptx.util.Pt = _Pt
    pptx.exc.PackageNotFoundError = _PackageNotFoundError
    pptx.Presentation = _Presentation

    sys.modules["pptx"] = pptx
    sys.modules["pptx.enum"] = pptx.enum
    sys.modules["pptx.enum.shapes"] = pptx.enum.shapes
    sys.modules["pptx.oxml"] = pptx.oxml
    sys.modules["pptx.oxml.ns"] = pptx.oxml.ns
    sys.modules["pptx.util"] = pptx.util
    sys.modules["pptx.exc"] = pptx.exc
