from collections import Counter
from pptx.enum.shapes import PP_PLACEHOLDER

from .base import Rule, Issue
from ..utils import iter_text_shapes


class FontRule(Rule):
    category = "Шрифты"

    def __init__(self, max_fonts: int = 3):
        self.max_fonts = max_fonts

    def run(self, pres) -> list[Issue]:
        issues: list[Issue] = []

        all_fonts: set[str] = set()
        title_ref_font: str | None = None

        for slide_idx, slide in enumerate(pres.slides, start=1):
            slide_fonts = self._fonts_on_slide(slide)
            all_fonts |= slide_fonts

            if len(slide_fonts) > 1:
                issues.append(
                    self._issue(
                        slide=slide_idx,
                        message=f"разные шрифты на слайде ({', '.join(sorted(slide_fonts))}).",
                    )
                )

            title_shape = self._get_title_shape(slide)
            if title_shape is not None:
                title_font = self._primary_font(title_shape)
                if title_font:
                    if title_ref_font is None:
                        title_ref_font = title_font
                    elif title_font != title_ref_font:
                        issues.append(
                            self._issue(
                                slide=slide_idx,
                                message=(
                                    f"заголовок не совпадает по шрифту с другими заголовками "
                                    f"({title_font} vs {title_ref_font})."
                                ),
                            )
                        )

        if self.max_fonts is not None and len(all_fonts) > self.max_fonts:
            issues.append(
                self._issue(
                    slide=None,
                    message=(
                        f"слишком много разных шрифтов во всей презентации "
                        f"({len(all_fonts)} > {self.max_fonts}): {', '.join(sorted(all_fonts))}."
                    ),
                )
            )

        return issues

    def _fonts_on_slide(self, slide) -> set[str]:
        fonts: set[str] = set()
        for shape in iter_text_shapes(slide):
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for p in getattr(tf, "paragraphs", []) or []:
                p_font = getattr(getattr(p, "font", None), "name", None)
                runs = getattr(p, "runs", []) or []
                if runs:
                    for r in runs:
                        r_font = getattr(getattr(r, "font", None), "name", None)
                        name = r_font or p_font
                        if name:
                            fonts.add(name)
                else:
                    if p_font:
                        fonts.add(p_font)
        return fonts

    def _get_title_shape(self, slide):
        try:
            title = slide.shapes.title
            if title is not None:
                return title
        except Exception:
            pass

        for shape in getattr(slide, "shapes", []) or []:
            if self._is_title(shape):
                return shape
        return None

    def _is_title(self, shape) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False

        if getattr(shape, "is_placeholder", False):
            try:
                ph_type = shape.placeholder_format.type
                return ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            except Exception:
                pass

        name = (getattr(shape, "name", "") or "").lower()
        if "title" in name:
            return True

        return False

    def _primary_font(self, shape) -> str | None:
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            return None

        fonts: list[str] = []
        for p in getattr(tf, "paragraphs", []) or []:
            p_font = getattr(getattr(p, "font", None), "name", None)
            runs = getattr(p, "runs", []) or []
            if runs:
                for r in runs:
                    r_font = getattr(getattr(r, "font", None), "name", None)
                    name = r_font or p_font
                    if name:
                        fonts.append(name)
            else:
                if p_font:
                    fonts.append(p_font)

        if not fonts:
            return None
        return Counter(fonts).most_common(1)[0][0]
