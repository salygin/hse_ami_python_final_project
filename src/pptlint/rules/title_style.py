from collections import Counter
from collections.abc import Iterable
from dataclasses import dataclass

from pptx.enum.shapes import PP_PLACEHOLDER

from .base import Rule, Issue


@dataclass(frozen=True)
class TitleSignature:
    font_name: str | None
    font_size_pt: float | None
    bold: bool | None


class TitleStyleRule(Rule):
    category = "Стиль заголовков"

    def __init__(self, *, ignore_first_slide: bool = True) -> None:
        self.ignore_first_slide = ignore_first_slide

    def run(self, pres) -> list[Issue]:
        issues: list[Issue] = []

        title_sigs: list[tuple[int, TitleSignature]] = []

        slide_height_emu = int(getattr(pres, "slide_height", 0) or 0)

        for idx, slide in enumerate(pres.slides):
            slide_no = idx + 1
            if self.ignore_first_slide and idx == 0:
                continue

            title_shape = self._find_title_shape(slide, slide_height_emu)
            if title_shape is None or not getattr(title_shape, "has_text_frame", False):
                issues.append(self._issue("на слайде нет заголовка.", slide=slide_no))
                continue

            title_text = (title_shape.text_frame.text or "").strip()
            if not title_text:
                parts = []
                for p in getattr(title_shape.text_frame, "paragraphs", []) or []:
                    part = (getattr(p, "text", "") or "").strip()
                    if part:
                        parts.append(part)
                title_text = " ".join(parts).strip()
            if not title_text:
                issues.append(self._issue("заголовок пустой.", slide=slide_no))
                continue

            sig = self._title_signature(title_shape)
            title_sigs.append((slide_no, sig))

        if not title_sigs:
            return issues

        baseline = Counter(sig for _, sig in title_sigs).most_common(1)[0][0]

        for slide_no, sig in title_sigs:
            if not self._signatures_equal(baseline, sig):
                issues.append(
                    self._issue(
                        "заголовок отличается по стилю от остальных "
                        f"(ожидается: {self._sig_to_str(baseline)}; найдено: {self._sig_to_str(sig)}).",
                        slide=slide_no,
                    )
                )

        return issues

    def _find_title_shape(self, slide, slide_height_emu: int):
        try:
            t = slide.shapes.title
            if t is not None and getattr(t, "has_text_frame", False):
                return t
        except Exception:
            pass

        for shape in getattr(slide, "shapes", []):
            if not getattr(shape, "has_text_frame", False):
                continue
            if self._is_title_placeholder(shape):
                return shape

        candidates = []
        top_limit = int(slide_height_emu * 0.28) if slide_height_emu else None

        for shape in getattr(slide, "shapes", []):
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = (tf.text or "").strip()
            if not text:
                continue

            if len(text) > 180:
                continue
            if len(getattr(tf, "paragraphs", [])) > 3:
                continue

            top_emu = int(getattr(shape, "top", 0) or 0)
            if top_limit is not None and top_emu > top_limit:
                continue

            max_size = self._max_font_size_pt(tf)
            score = (max_size or 0.0) * 10.0 - (top_emu / 100000.0)
            candidates.append((score, shape))

        if not candidates:
            return None

        candidates.sort(key=lambda x: x[0], reverse=True)
        return candidates[0][1]

    def _is_title_placeholder(self, shape) -> bool:
        try:
            if not getattr(shape, "is_placeholder", False):
                return False
            ph_type = shape.placeholder_format.type
            return ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        except Exception:
            return False


    def _title_signature(self, shape) -> TitleSignature:
        tf = shape.text_frame

        font_names: list[str] = []
        font_sizes: list[float] = []
        bold_vals: list[bool] = []

        for run in self._iter_runs(tf):
            f = getattr(run, "font", None)
            if f is None:
                continue

            name = getattr(f, "name", None)
            if isinstance(name, str) and name.strip():
                font_names.append(name.strip())

            size = getattr(f, "size", None)
            if size is not None:
                try:
                    font_sizes.append(float(size.pt))
                except Exception:
                    pass

            b = getattr(f, "bold", None)
            if b is True:
                bold_vals.append(True)
            elif b is False:
                bold_vals.append(False)

        font_name = self._mode_str(font_names)
        font_size_pt = self._mode_float(font_sizes, round_to=0.5)
        bold = self._mode_bool(bold_vals)

        return TitleSignature(font_name=font_name, font_size_pt=font_size_pt, bold=bold)

    def _iter_runs(self, text_frame) -> Iterable[object]:
        for p in getattr(text_frame, "paragraphs", []):
            for r in getattr(p, "runs", []):
                yield r

    def _max_font_size_pt(self, text_frame) -> float | None:
        mx: float | None = None
        for run in self._iter_runs(text_frame):
            f = getattr(run, "font", None)
            if f is None:
                continue
            size = getattr(f, "size", None)
            if size is None:
                continue
            try:
                pt = float(size.pt)
            except Exception:
                continue
            mx = pt if mx is None else max(mx, pt)
        return mx


    def _signatures_equal(self, a: TitleSignature, b: TitleSignature) -> bool:
        if a.font_name is not None and b.font_name is not None and a.font_name != b.font_name:
            return False
        if a.font_size_pt is not None and b.font_size_pt is not None and abs(a.font_size_pt - b.font_size_pt) > 0.01:
            return False
        if a.bold is not None and b.bold is not None and a.bold != b.bold:
            return False
        return True

    def _sig_to_str(self, sig: TitleSignature) -> str:
        name = sig.font_name or "?"
        size = f"{sig.font_size_pt:g}pt" if sig.font_size_pt is not None else "?pt"
        bold = "bold" if sig.bold is True else ("regular" if sig.bold is False else "?")
        return f"{name}, {size}, {bold}"

    def _mode_str(self, xs: list[str]) -> str | None:
        if not xs:
            return None
        return Counter(xs).most_common(1)[0][0]

    def _mode_float(self, xs: list[float], *, round_to: float = 1.0) -> float | None:
        if not xs:
            return None
        if round_to > 0:
            xs = [round(x / round_to) * round_to for x in xs]
        return Counter(xs).most_common(1)[0][0]

    def _mode_bool(self, xs: list[bool]) -> bool | None:
        if not xs:
            return None
        return Counter(xs).most_common(1)[0][0]
