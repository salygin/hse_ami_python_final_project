from pptx.oxml.ns import qn

from .base import Rule, Issue
from ..utils import iter_text_shapes


class ListRule(Rule):
    category = "Списки"

    def run(self, pres) -> list[Issue]:
        issues: list[Issue] = []

        for slide_idx, slide in enumerate(getattr(pres, "slides", []), start=1):
            list_frames = 0
            one_item_frames = 0
            mixed_frames = 0

            for shape in iter_text_shapes(slide):
                tf = shape.text_frame
                items = self._count_items(tf)
                if items <= 0:
                    continue

                list_frames += 1
                if items == 1:
                    one_item_frames += 1
                if self._has_mixed_bullets(tf):
                    mixed_frames += 1

            if one_item_frames > 0:
                if list_frames >= 2 and one_item_frames == 1:
                    msg = f"Слайд {slide_idx}: на слайде {list_frames} списка, один из них из 1 пункта."
                elif list_frames == 1 and one_item_frames == 1:
                    msg = f"Слайд {slide_idx}: список из 1 пункта."
                else:
                    msg = f"Слайд {slide_idx}: {one_item_frames} списка из 1 пункта."
                issues.append(self._issue(msg, slide=slide_idx))

            if mixed_frames > 0:
                if mixed_frames == 1:
                    msg = (
                        f"Слайд {slide_idx}: список с перемешанными маркерами."
                    )
                else:
                    msg = (
                        f"Слайд {slide_idx}: {mixed_frames} списка с перемешанными маркерами."
                    )
                issues.append(self._issue(msg, slide=slide_idx))

        return issues

    def _count_items(self, text_frame) -> int:
        cnt = 0
        for p in getattr(text_frame, "paragraphs", []):
            text = (getattr(p, "text", "")).strip()
            if not text:
                continue
            if self._is_list_paragraph(p):
                cnt += 1
        return cnt

    def _has_mixed_bullets(self, text_frame) -> bool:
        types: list[str] = []
        for p in getattr(text_frame, "paragraphs", []):
            text = (getattr(p, "text", "")).strip()
            if not text:
                continue
            types.append(self._paragraph_type(p))

        if len(types) < 2:
            return False

        transitions = sum(1 for i in range(1, len(types)) if types[i] != types[i - 1])
        return transitions >= 2

    @staticmethod
    def _paragraph_type(paragraph) -> str:
        p = getattr(paragraph, "_p", None)
        if p is None:
            return "none"

        pPr = getattr(p, "pPr", None)
        if pPr is None:
            return "none"

        if pPr.find(qn("a:buNone")) is not None:
            return "none"

        if pPr.find(qn("a:buAutoNum")) is not None:
            return "auto"
        if pPr.find(qn("a:buChar")) is not None:
            return "char"
        if pPr.find(qn("a:buBlip")) is not None:
            return "blip"
        if pPr.find(qn("a:buFont")) is not None:
            return "font"

        return False
