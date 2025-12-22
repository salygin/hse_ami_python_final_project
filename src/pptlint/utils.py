from collections.abc import Iterator

from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_text_shapes(slide: object) -> Iterator[object]:
    yield from _walk_shapes(getattr(slide, "shapes", []))


def _walk_shapes(shapes: object) -> Iterator[object]:
    for shape in shapes:
        if _is_group_shape(shape):
            inner = getattr(shape, "shapes", None)
            if inner is not None:
                yield from _walk_shapes(inner)
            continue

        if getattr(shape, "has_text_frame", False):
            yield shape


def _is_group_shape(shape: object) -> bool:
    st = getattr(shape, "shape_type", None)
    if st == MSO_SHAPE_TYPE.GROUP:
        return True
    return hasattr(shape, "shapes") and not getattr(shape, "has_text_frame", False)
